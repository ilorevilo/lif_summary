# -*- coding: utf-8 -*-
"""
    lif_summary: easy creation and extraction of microscope images/ videos from a .lif-file
    written by Oliver Schneider (08/2017)
  
"""

import sys

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import numpy as np
import javabridge
import bioformats
from xml.etree import ElementTree as ETree
from xml import etree as et

from skimage.viewer import ImageViewer
from skimage import exposure
import skimage
import PIL
from PIL import ImageDraw, ImageFont, Image

import cv2

from itertools import zip_longest

import glob
import os
import tifffile

class lif_summary:
    """
        class which contains all functions to extract images from lif file
    """

    def __init__(self, lif_file):
        """
            just specify path to lif-file in constructor
        """
        self.lif_file = lif_file
        self.filename = os.path.splitext(self.lif_file)[0]
        self.reader = bioformats.ImageReader(self.lif_file, perform_init=True)
        self.categorized_series = {"img_simple":[],"img_multiC":[],"img_multiT":[],"img_multiTC":[],"other":[],"z_stack":[]}
        # specify categories under which images "series" will get sorted later
    
    def get_hist_boundaries(self, omeobject, image_name):
        """
            extracts histogram boundaries, stored in omeobject
        """
        rootnode = omeobject.dom.getroot()
        root_iter = rootnode.iter()

        white_val, black_val = 0, 0

        for child in root_iter:
            if child.text == image_name + ' Image|ChannelScalingInfo|WhiteValue':
                # white_val = root_iter.next().text
                white_val = next(root_iter).text    #changed for Python 3
                break

        root_iter = rootnode.iter()

        for child in root_iter:
            if child.text == image_name + ' Image|ChannelScalingInfo|BlackValue':
                black_val = next(root_iter).text
                break
        return black_val, white_val	

    def get_bit_resolution(self, omeobject, image_name):
        """
            extracts bit resolution, stored in omeobject
        """
        rootnode = omeobject.dom.getroot()
        root_iter = rootnode.iter()

        resolution = 0

        for child in root_iter:
            if child.text == image_name + ' Image|ChannelDescription|Resolution':
                resolution = next(root_iter).text
                break

        return resolution
        
    def get_image_overview(self):
        """
            fills categorized_series dict with correct entries from lif-file
            basically: connect image name/type to id
        """
            
        metadata = bioformats.get_omexml_metadata(self.lif_file)
        metadata = metadata.encode('utf8')
        metadata = metadata.decode('utf8')    #here decode for string...    #needed?
        ome = bioformats.OMEXML(metadata)
        
        image_count = ome.image_count
        print ("------------------------------")
        print ("reading file", self.filename)
        print ("number of image-series in lif-file: " + str(image_count))
            
        for image_nr in range(image_count):
            iome = ome.image(image_nr)    #iome: current image, loop through all images
            
            """ tries to extract metadata
            #origmeta = {}
            #ome_origmeta= ome.OriginalMetadata(origmeta)#does not need dict, class is inherited from dict
            #print(ome_origmeta)
            
            #structannot = ome.StructuredAnnotations    #(image_nr) # use property...
            #not working?
            
            #plane = ome.plane()
            #print(iome.Pixels.Plane().DeltaT)   #specify nr of plane

            #structannot.keys()
            #structannot.iter_original_metadata()
            """
            
            # create dict with information of current series
            new_imagentry = {}
                    
            new_imagentry['name'] = iome.get_Name()
            new_imagentry['id'] = iome.get_ID()
            new_imagentry['nr'] = image_nr
            new_imagentry['mpp'] = iome.Pixels.PhysicalSizeX # is returned in mpp
            
            new_imagentry['SizeX'] = iome.Pixels.get_SizeX()
            new_imagentry['SizeT'] = iome.Pixels.get_SizeT()
            new_imagentry['SizeZ'] = iome.Pixels.get_SizeZ()
            new_imagentry['SizeC'] = iome.Pixels.get_SizeC()
            
            black_val, white_val = self.get_hist_boundaries(ome,iome.get_Name())
            new_imagentry['Blackval'] = float(black_val)
            new_imagentry['Whiteval'] = float(white_val)
            new_imagentry['bit_resolution'] = int(self.get_bit_resolution(ome, iome.get_Name()))
            
            print ("hist_ranges:")
            print (black_val, white_val)
            
            new_imagentry['subfolder'] = None   #specify if images (like mark and find) are in a subfolder. currently only 1 level supported
            
            # calculate fps for videos, assumes evenly spaced frames
            if new_imagentry['SizeT'] > 1:
                last_plane = new_imagentry['SizeT']-1
                new_imagentry['duration'] = iome.Pixels.Plane(last_plane).DeltaT
                new_imagentry['fps'] = last_plane/new_imagentry['duration'] #calculate fps
            
            #print(new_imagentry)
            ## sort dict according to entries
            
            # move z-stack to category z_stack
            # todo in future: check for multichannel z-stacks
            if new_imagentry['SizeZ'] > 1:
                self.categorized_series["z_stack"].append(new_imagentry) 
            
            # move EnvironmentalGraph-files into category other
            elif "EnvironmentalGraph" in new_imagentry["name"]:
                self.categorized_series["other"].append(new_imagentry)
            
            # video
            elif new_imagentry['SizeT'] > 1:
                
                #multichannel video
                if new_imagentry['SizeC'] > 1:
                    self.categorized_series["img_multiTC"].append(new_imagentry)
                
                #normal video
                elif new_imagentry['SizeC'] == 1:
                    
                    if "Mark_and_Find" in new_imagentry["name"]:
                        new_imagentry['subfolder'] = new_imagentry['name'].split("/")[0]# save folder name to create later
                        new_imagentry['name'] = new_imagentry['name'].split("/")[1]
                    else:
                        new_imagentry['name'] = new_imagentry['name'].split("/")[0] #can cause issues if in subfolder
                    self.categorized_series["img_multiT"].append(new_imagentry)
                
                #everything else
                else:
                    self.categorized_series["other"].append(new_imagentry)
                
            # no time variation
            elif new_imagentry['SizeT'] == 1:
                
                #multichannel image
                if new_imagentry['SizeC'] > 1:
                    if "EnvironmentalGraph" not in new_imagentry["name"]:
                        self.categorized_series["img_multiC"].append(new_imagentry)
                
                #normal image
                elif new_imagentry['SizeC'] == 1:
                    
                    if "Mark_and_Find" in new_imagentry["name"]:
                        new_imagentry['subfolder'] = new_imagentry['name'].split("/")[0]# save folder name to create later
                        #new_imagentry['name'] = new_imagentry['name'].split("/")[1]
                    else:
                        new_imagentry['name'] = new_imagentry['name'].split("/")[0] #can cause issues if in subfolder
                    self.categorized_series["img_simple"].append(new_imagentry)
                
                #everything else
                else:
                    self.categorized_series["other"].append(new_imagentry)  
            
            else:
                self.categorized_series["other"].append(new_imagentry)
                
    def create_output_folders(self):
        """
            creates output folder structure based on filename
        """
        
        outputfolder = self.filename

        #create folder with name of lif_file if it doesn't exist
        if os.path.isdir(outputfolder):
            print ("folder " + outputfolder + " already exists")
            print ("don't create folder structure")
            print ("don't export folders, quitting...")
            
            return False
        else:
            print ("create folder " + outputfolder)
            print ("------------------------------")
            os.makedirs(outputfolder)

            os.makedirs(os.path.join(outputfolder,"raw tifs","images"))
            os.makedirs(os.path.join(outputfolder,"raw tifs","videos"))
            os.makedirs(os.path.join(outputfolder,"raw tifs","multichannel images"))
            os.makedirs(os.path.join(outputfolder,"raw tifs","zstacks")) # do it smarter in future and just create folders if imagetypes exist...
            os.makedirs(os.path.join(outputfolder,"compressed","images"))
            os.makedirs(os.path.join(outputfolder,"compressed","videos"))
            #os.makedirs(os.path.join(outputfolder,"compressed","multichannel images")) #does not work yet
            
            return True
            
    def save_single_tif(self, image, path, resolution_mpp, photometric = None):
        """
            saves single tif into specified folder with specified meta
        """
        resolution_ppm = 1/resolution_mpp# convert micron per pixel to pixel per micron
        if photometric == None:
            tifffile.imsave(path, image, imagej=True, resolution=(resolution_ppm, resolution_ppm), metadata={'unit': 'um'}) #what if ppm is different in x and y? #0.6149
        else:
            tifffile.imsave(path, image, imagej=True, resolution=(resolution_ppm, resolution_ppm), metadata={'unit': 'um'}, photometric = photometric)
        

    def get_image_array(self, c, z, t, series):
        """
            reads image via bioformats reader to np arary
        """
        image = self.reader.read(c=c, z=z, t=t, series=series, rescale=False)
        
        image[image < 0] = image[image < 0] + 65535 #weird results were happening here when 12 bit data was saved in 16 bit container?
        image_uint = np.uint16(image)
        return image_uint

    def export_raw_tifs(self):
        """
            exports raw data saved in lif-file
        """
        
        print("exporting raw tifs")
        
        # export simple images
        total_images = len(self.categorized_series["img_simple"])
        print("exporting " + str(total_images) + " simple images")

        for nr, imagentry in enumerate(self.categorized_series["img_simple"]):
            
            self.printProgressBar(nr + 1, total_images, prefix = 'Progress:', suffix = 'Complete', length = 50)

            series_nr = imagentry["nr"]
            series_name = imagentry["name"]
            resolution_mpp = imagentry["mpp"]
            outputpath = os.path.join(self.filename,"raw tifs","images",series_name+".tif")
            if imagentry["subfolder"] != None:
                subfolder = os.path.join(self.filename,"raw tifs","images",imagentry["subfolder"])
                if not os.path.exists(subfolder):
                    os.makedirs(subfolder)
            else:
                subfolder = None
            
            image = self.get_image_array(c=0, z=0, t=0, series = series_nr)
            
            # create output folder if MarkAndFind?

            self.save_single_tif(image, outputpath, resolution_mpp)

        # export videos
        total_videos = len(self.categorized_series["img_multiT"])
        print("exporting " + str(total_videos) + " videos")
        
        for videonr, imagentry in enumerate(self.categorized_series["img_multiT"]):
            series_name = imagentry["name"]
            print("video", videonr)
            print(series_name)
            
            subfolder = imagentry["subfolder"]

            # create subfolder, depending on MarkAndFind
            if subfolder != None:
                outputpath = os.path.join(self.filename,"raw tifs","videos", subfolder, series_name)
            else:
                outputpath = os.path.join(self.filename,"raw tifs","videos", series_name)
            
            if not os.path.exists(outputpath):
                    os.makedirs(outputpath)
            
            series_nr = imagentry["nr"]            
            resolution_mpp = imagentry["mpp"]
            totalframes = imagentry['SizeT']
            
            for frame in np.arange(totalframes):
                self.printProgressBar(frame+1, totalframes, prefix = 'Progress:', suffix = 'Complete', length = 50)
                
                framepath = os.path.join(outputpath, series_name+"-{:04d}.tif".format(frame))
                image = self.get_image_array(c=0, z=0, t=frame, series = series_nr)
                self.save_single_tif(image,framepath, resolution_mpp)            
            
            bit_resolution = imagentry["bit_resolution"]
            image_scale = 2**bit_resolution - 1
            
            # write dict with metainfos into videoinfos.txt
            metafile = os.path.join(outputpath, "videoinfos.txt")
            metadict = {'microns_per_pixel' : resolution_mpp, 'fps' : imagentry["fps"], 'Blackval':imagentry["Blackval"]*image_scale, 'Whiteval':imagentry["Whiteval"]*image_scale}
            f = open(metafile,"w")
            f.write( str(metadict) )
            f.close()
        
        # export multichannel images
        total_multiC = len(self.categorized_series["img_multiC"])
        print("exporting " + str(total_multiC) + " multichannel images")       
        
        for nr, imagentry in enumerate(self.categorized_series["img_multiC"]):
            
            self.printProgressBar(nr + 1, total_multiC, prefix = 'Progress:', suffix = 'Complete', length = 50)
            
            series_nr = imagentry["nr"]
            series_name = imagentry["name"]
            resolution_mpp = imagentry["mpp"]
            totalchannels = imagentry['SizeC']
            outputpath = os.path.join(self.filename,"raw tifs","multichannel images",series_name+".tif")
            
            image = self.get_image_array(c= None, z=0, t=0, series = series_nr)
            image = np.moveaxis(image,2,0)  #exchange dimensions
            self.save_single_tif(image,outputpath, resolution_mpp, photometric = 'minisblack')      
    
        # export zstacks images
        total_zstack = len(self.categorized_series["z_stack"])
        print("exporting " + str(total_zstack) + " zstacks") 
        
        for nr, imagentry in enumerate(self.categorized_series["z_stack"]):

            series_nr = imagentry["nr"]
            series_name = imagentry["name"]
            resolution_mpp = imagentry["mpp"]
            planes = imagentry['SizeZ']
            #totalchannels = imagentry['SizeC']
            print("exporting z-stack ", series_name)
            
            outputpath = os.path.join(self.filename,"raw tifs","zstacks",series_name)  

            if not os.path.exists(outputpath):
                    os.makedirs(outputpath)

            for plane in np.arange(planes):
                self.printProgressBar(plane+1, planes, prefix = 'Progress:', suffix = 'Complete', length = 50)
                
                #outputpath = os.path.join(self.filename,"raw tifs","videos", series_name, series_name+"-{:04d}.tif".format(frame))            
                planepath = os.path.join(outputpath, series_name+"-{:04d}.tif".format(plane))
                image = self.get_image_array(c=0, z=plane, t=0, series = series_nr)
                self.save_single_tif(image,planepath, resolution_mpp)            
    
    def adj_contrast(self, image, vmin=None, vmax=None):
        """
            adusts contrast of inputimage either by setting to specified values or by calculating the 98% range
        """
        if None in (vmin, vmax):
            vmin = np.percentile(image, 0.2)
            vmax = np.percentile(image, 99.8)
                
        #print("adjusting contrast to: ", vmin, vmax)
        img_adj_contrast = exposure.rescale_intensity(image, in_range=(vmin, vmax))
        
        return img_adj_contrast

    def create_scalebar(self, dimX_px, microns_per_pixel):
        """
            creates scalebar as np array which then can be transferred to image
        """
        
        scale_values = [1,2,5,10,15,20,30,40,50,70,100,150,200,300,400,500,700,1000]
        initial_scale_length = dimX_px * 0.2 * microns_per_pixel
        
        text_height_px = int(round(dimX_px * 0.05))
        drawfont = ImageFont.truetype("arial.ttf", text_height_px)
        
        scale_length_microns = min(scale_values, key=lambda x:abs(x-initial_scale_length))    # pick nearest value
        scale_caption = str(scale_length_microns) + " µm"
        scale_length_px = scale_length_microns / microns_per_pixel
        scale_height_px = dimX_px * 0.01
        
        bg_square_spacer_px = scale_length_px * 0.07
        
        bg_square_length_px = int(round(scale_length_px + 2 * bg_square_spacer_px))
        bg_square_height_px = int(round(text_height_px + scale_height_px + 2 * bg_square_spacer_px))
        
        scalebar = Image.new("L", (bg_square_length_px, bg_square_height_px), "white")
        draw = ImageDraw.Draw(scalebar)
        
        w_caption, h_caption = draw.textsize(scale_caption, font = drawfont)        
        
        draw.rectangle(((0, 0), (bg_square_length_px, bg_square_height_px)), fill="black")
        draw.rectangle(((bg_square_spacer_px, bg_square_height_px - bg_square_spacer_px - scale_height_px), (bg_square_length_px - bg_square_spacer_px, bg_square_height_px - bg_square_spacer_px)), fill="white")
        draw.text((bg_square_spacer_px + bg_square_length_px/2 - w_caption/2, bg_square_spacer_px/2), scale_caption, font = drawfont, fill="white")
        
        output_scalebar = np.array(scalebar)
        return output_scalebar

    def plot_scalebar(self, input_image, microns_per_pixel, image_name = None):
        """
            plots scalebar + title onto image if desired
            image input: np array
        """
        
        image_scalebar = PIL.Image.fromarray(input_image)   #np.uint8(input_image*255)
        
        dimX_px = input_image.shape[0]
        initial_scale_length = dimX_px * 0.2 * microns_per_pixel
        
        text_height_px = int(round(dimX_px * 0.05))
        
        scale_values = [1,2,5,10,15,20,30,40,50,70,100,150,200,300,400,500,700,1000]
        drawfont = ImageFont.truetype("arial.ttf", text_height_px)
        
        scale_length_microns = min(scale_values, key=lambda x:abs(x-initial_scale_length))    # pick nearest value
        scale_caption = str(scale_length_microns) + " µm"
        
        draw = ImageDraw.Draw(image_scalebar)
        w_caption, h_caption = draw.textsize(scale_caption, font = drawfont)
        
        scale_length_px = scale_length_microns / microns_per_pixel
        scale_height_px = dimX_px * 0.01
        
        bg_square_spacer_px = scale_length_px * 0.07
        
        bg_square_length_px = scale_length_px + 2 * bg_square_spacer_px
        bg_square_height_px = text_height_px + scale_height_px + 2 * bg_square_spacer_px
                
        if dimX_px > 800:
        
            draw.rectangle(((dimX_px - bg_square_length_px, dimX_px - bg_square_height_px), (dimX_px, dimX_px)), fill="black")
            draw.rectangle(((dimX_px - bg_square_length_px + bg_square_spacer_px, dimX_px - bg_square_spacer_px - scale_height_px), (dimX_px - bg_square_spacer_px, dimX_px - bg_square_spacer_px)), fill="white")
            draw.text((dimX_px - bg_square_length_px + bg_square_spacer_px + bg_square_length_px/2 - w_caption/2, dimX_px - bg_square_height_px + bg_square_spacer_px/2), scale_caption, font = drawfont, fill="white")# scale_caption.decode('utf8')

            # burn title if provided
            if image_name != None:
                title_height_px = int(round(dimX_px * 0.05))
                drawfont = ImageFont.truetype("arial.ttf", title_height_px)
                draw.rectangle(((0,0),(dimX_px,title_height_px*1.2)), fill = "black")
                draw.text((0,0),image_name, font = drawfont, fill = "white")
        
        output_image = np.array(image_scalebar)
                
        return output_image

    def printProgressBar (self, iteration, total, prefix = '', suffix = '', decimals = 1, length = 100, fill = '█'):
        """
        print iteration progress
        Call in a loop to create terminal progress bar
        @params:
            iteration   - Required  : current iteration (Int)
            total       - Required  : total iterations (Int)
            prefix      - Optional  : prefix string (Str)
            suffix      - Optional  : suffix string (Str)
            decimals    - Optional  : positive number of decimals in percent complete (Int)
            length      - Optional  : character length of bar (Int)
            fill        - Optional  : bar fill character (Str)
        """
        percent = ("{0:." + str(decimals) + "f}").format(100 * (iteration / float(total)))
        filledLength = int(length * iteration // total)
        bar = fill * filledLength + '-' * (length - filledLength)
        print('\r%s |%s| %s%% %s' % (prefix, bar, percent, suffix), end = '\r')
        # Print New Line on Complete
        if iteration == total: 
            print()
    
    def export_compressed_data(self):
        """
            exports jpgs and videos of files, adds title + scale bar for jpgs, scale bar for video
        """
        
        print("-----------------------------")
        print("exporting compressed data")
        
        #make it easier, just one loop and separare then between raw and compressed?
        
        #images
        total_images = len(self.categorized_series["img_simple"])
        print("exporting " + str(total_images) + " simple images")

        for nr, imagentry in enumerate(self.categorized_series["img_simple"]):

            self.printProgressBar(nr + 1, total_images, prefix = 'Progress:', suffix = 'Complete', length = 50)
            
            series_nr = imagentry["nr"]
            series_name = imagentry["name"]
            resolution_mpp = imagentry["mpp"]
            outputpath_jpg = os.path.join(self.filename,"compressed","images",series_name+".jpg")
            
            if imagentry["subfolder"] != None:
                subfolder = os.path.join(self.filename,"compressed","images",imagentry["subfolder"])
                if not os.path.exists(subfolder):
                    os.makedirs(subfolder)
            else:
                subfolder = None            
            
            image = self.get_image_array(c=0, z=0, t=0, series = series_nr)
            bit_resolution = imagentry["bit_resolution"]
            image_scale = 2**bit_resolution - 1
            
            #image_adj_contrast = self.adj_contrast(image, imagentry["Blackval"], imagentry["Whiteval"])
            image_adj_contrast = self.adj_contrast(image, imagentry["Blackval"]*image_scale, imagentry["Whiteval"]*image_scale)
            image_8 = cv2.convertScaleAbs(image_adj_contrast,alpha=(255.0/65535.0))
            
            labeled_image = self.plot_scalebar(image_8, resolution_mpp, series_name)
            
            skimage.io.imsave(outputpath_jpg, labeled_image)
            
        # videos
        
        total_videos = len(self.categorized_series["img_multiT"])
        print("exporting " + str(total_videos) + " videos")
        
        fourcc = cv2.VideoWriter_fourcc(*'MP4V')
        
        for videonr, imagentry in enumerate(self.categorized_series["img_multiT"]):
   	
            print("video", videonr)
            
            series_name = imagentry["name"]
            series_nr = imagentry["nr"]            
            resolution_mpp = imagentry["mpp"]
            totalframes = imagentry['SizeT']
            sizeX = imagentry['SizeX']
            fps = imagentry['fps']
            
            if fps < 0.1:   #timelapse video
                fps = 5
            
            if sizeX == 2048:   #scale down 2048 videos only
                scalefactor = 2
            else:
                scalefactor = 1
            outputsizeX = int(sizeX/scalefactor)
            
            subfolder = imagentry["subfolder"]

            # create subfolder, depending on MarkAndFind
            if subfolder != None:
                outputfolder = os.path.join(self.filename,"compressed","videos", subfolder)
            else:
                outputfolder = os.path.join(self.filename,"compressed","videos")
            
            if not os.path.exists(outputfolder):
                os.makedirs(outputfolder)            
            
            
            outputpath = os.path.join(outputfolder, series_name + ".mp4")
            video=cv2.VideoWriter(outputpath,fourcc,fps,(outputsizeX,outputsizeX),False)#-1: Select codec manually, False: monochrome images

            print(outputpath)

            scalebar = self.create_scalebar(sizeX,resolution_mpp)# don't adjust, applied to raw image, scaled down together
            scale_width_px = scalebar.shape[0]
            scale_height_px = scalebar.shape[1]
            
            for frame in np.arange(totalframes):
                
                self.printProgressBar(frame+1, totalframes, prefix = 'Progress:', suffix = 'Complete', length = 50)
                
                image = self.get_image_array(c=0, z=0, t=frame, series = series_nr)
                bit_resolution = imagentry["bit_resolution"]
                image_scale = 2**bit_resolution - 1
                
                image_adj_contrast = self.adj_contrast(image, imagentry["Blackval"]*image_scale, imagentry["Whiteval"]*image_scale)   # each image gets adjusted individually... do once and keep values for all
                image_8 = cv2.convertScaleAbs(image_adj_contrast,alpha=(255.0/65535.0))                              
                image_8[-1-scale_width_px:-1,-1-scale_height_px:-1] = scalebar
                image_8 = cv2.resize(image_8,(outputsizeX,outputsizeX))
                
                video.write(image_8)
            cv2.destroyAllWindows()
            video.release() 

    def create_ppt_summary(self):
        """
            creates ppt with all exported images
        """
        
        print("creating ppt-summary")
        
        # take always n elements from list as list, helpfunction
        def grouper(iterable, n, fillvalue=None):
            args = [iter(iterable)] * n
            return zip_longest(*args, fillvalue=fillvalue)
        
        # prepare presentation
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = str(self.filename)
        subtitle.text = "lif-summary"

        ###############################
        # constants for 2 x 3 layout

        Pt_per_cm = 72.0/2.54

        slide_width = Pt (25.4 * Pt_per_cm)
        slide_height = Pt (19.05 * Pt_per_cm)

        headspace = Pt(30)
        image_height = Pt(200)

        d_horizontal = (slide_width-3*image_height)/4
        d_vertical = (slide_height-headspace-2*image_height)/3

        ##############################
        
        # pick 6 images
        
        for images_6group in grouper(self.categorized_series['img_simple'], 6, None):

            #add new slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            #iterate through rows, columns
            for rowindex in range(2):
                for columnindex in range(3):
                    
                    image_index = rowindex * 3 + columnindex
                    if(images_6group[image_index]) != None:
                        image_name = images_6group[image_index]['name']
                        imagepath = os.path.join(self.filename,"compressed","images",image_name + ".jpg")
                        pic = slide.shapes.add_picture(imagepath, d_horizontal * (columnindex + 1) + image_height * columnindex, headspace + d_vertical * (rowindex + 1) + image_height * rowindex, height=image_height)
        
        """
        # pick videos, experimental, works but stillimage is loudspeaker -> create individual stillimage
        for video in self.categorized_series['img_multiT']:
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            videopath = os.path.join(self.filename,"compressed","videos",video['name'] + ".mp4")
            slide.shapes.add_movie(videopath,d_vertical,d_vertical,slide_height-2*d_vertical,slide_height-2*d_vertical)
        """
        
        outputpath = os.path.join(str(self.filename),str(self.filename) + '_summary.pptx')
        prs.save(outputpath)
     
    def close(self):
        """
            closes bioformats reader
        """

        self.reader.close()
        bioformats.release_image_reader(self.lif_file)
        bioformats.clear_image_reader_cache()
        
def main():
    javabridge.start_vm(class_path=bioformats.JARS, max_heap_size='8G')
    # issues in ipython... start manually there
    
    inputlifs = glob.glob("*.lif") # get all .lif-files in current folder
    #inputlifs = [inputfile]    # if you just want to read one file
    
    for inputfile in inputlifs:     
        new_summary = lif_summary(inputfile)
        new_summary.get_image_overview()
        if new_summary.create_output_folders() == True:
            new_summary.export_raw_tifs()
            new_summary.export_compressed_data()
            new_summary.create_ppt_summary()
        new_summary.close()
    
    javabridge.kill_vm()
    
    # move files into newly created folders when finished
    for inputfile in inputlifs:
        print("moving lif-file")
        os.rename(inputfile, os.path.join(os.path.splitext(inputfile)[0], inputfile))
        print("finished")
    
                   
if __name__ == "__main__":
    main()